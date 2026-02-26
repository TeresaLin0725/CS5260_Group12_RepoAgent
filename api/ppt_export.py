"""
PPT Export Module — Presentation Renderer

Converts structured AnalyzedContent into a PowerPoint presentation.
Phase 2b-PPT (template assembly) + Phase 3 (python-pptx render).

Slide structure adapts to repo_type_hint:
  - Title slide (always)
  - Project Overview (always)
  - Architecture / Service Topology
  - Tech Stack
  - Key Modules / Components
  - Data Flow / Pipeline
  - API / Commands (emphasised for library / cli_tool)
  - Component Hierarchy (webapp only)
  - Deployment (microservice only)
  - Data Schemas (data_pipeline only)
  - Target Users / Closing (generic / data_pipeline)
"""

import logging
from io import BytesIO
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from api.content_analyzer import AnalyzedContent

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def _add_title_slide(prs, repo_name: str):
    """Slide 0: Title card."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(35, 55, 100)

    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = repo_name
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = 1  # center

    p2 = tf.add_paragraph()
    p2.text = "Architecture Overview"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(180, 200, 240)
    p2.alignment = 1


def _add_content_slide(prs, title: str, bullets: list[str], max_bullets: int = 6):
    """Generic content slide with title + bullet points."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(35, 55, 100)

    # Bullets
    txBody = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(5.5))
    tfb = txBody.text_frame
    tfb.word_wrap = True

    displayed = bullets[:max_bullets]
    for i, bullet in enumerate(displayed):
        if i == 0:
            para = tfb.paragraphs[0]
        else:
            para = tfb.add_paragraph()
        para.text = bullet
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(50, 50, 50)
        para.space_after = Pt(6)


def _add_paragraph_slide(prs, title: str, text: str):
    """Slide with a title and a paragraph (for overview, target users, etc.)."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(35, 55, 100)

    txBody = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(5.5))
    tfb = txBody.text_frame
    tfb.word_wrap = True
    para = tfb.paragraphs[0]
    para.text = text[:1200]  # safety limit
    para.font.size = Pt(14)
    para.font.color.rgb = RGBColor(50, 50, 50)


# ---------------------------------------------------------------------------
# Phase 2b-PPT: structured AnalyzedContent → slide outline
# ---------------------------------------------------------------------------

def _build_slides(prs, analyzed: "AnalyzedContent"):
    """Populate the presentation with slides adapted to repo_type_hint."""

    rtype = analyzed.repo_type_hint or "generic"

    # 1) Title
    _add_title_slide(prs, analyzed.repo_name)

    # 2) Overview (always)
    if analyzed.project_overview:
        _add_paragraph_slide(prs, "Project Overview", analyzed.project_overview)

    # 3) Architecture / Service Topology
    if analyzed.architecture:
        title = "Service Topology & Architecture" if rtype == "microservice" else "Architecture & Design"
        _add_content_slide(prs, title, analyzed.architecture)

    # 4) Tech Stack
    ts = analyzed.tech_stack
    if ts and (ts.languages or ts.frameworks or ts.key_libraries or ts.infrastructure):
        bullets = []
        if ts.languages:
            bullets.append("Languages: " + ", ".join(ts.languages))
        if ts.frameworks:
            bullets.append("Frameworks: " + ", ".join(ts.frameworks))
        if ts.key_libraries:
            bullets.append("Key Libraries: " + ", ".join(ts.key_libraries))
        if ts.infrastructure:
            bullets.append("Infrastructure: " + ", ".join(ts.infrastructure))
        _add_content_slide(prs, "Tech Stack", bullets)

    # 5) Key Modules
    if analyzed.key_modules:
        title = "Key Services" if rtype == "microservice" else "Key Modules & Components"
        bullets = [f"{m.name} — {m.responsibility}" for m in analyzed.key_modules]
        _add_content_slide(prs, title, bullets, max_bullets=7)

    # 6) Data Flow / Pipeline
    if analyzed.data_flow:
        title = "Data Flow & Pipeline Stages" if rtype == "data_pipeline" else "Data Flow & Processing"
        if rtype == "microservice":
            title = "Message Flow"
        _add_content_slide(prs, title, analyzed.data_flow)

    # 7) API / Commands (emphasised for library / cli_tool)
    if analyzed.api_points:
        title = "Commands & Interface" if rtype == "cli_tool" else "API & Integration Points"
        _add_content_slide(prs, title, analyzed.api_points)

    # 8) Optional: Component Hierarchy (webapp)
    if rtype == "webapp" and analyzed.component_hierarchy:
        _add_paragraph_slide(prs, "Component Hierarchy & Routing", analyzed.component_hierarchy)

    # 9) Optional: Deployment (microservice)
    if rtype == "microservice" and analyzed.deployment_info:
        _add_paragraph_slide(prs, "Deployment & Infrastructure", analyzed.deployment_info)

    # 10) Optional: Data Schemas (data_pipeline)
    if rtype == "data_pipeline" and analyzed.data_schemas:
        _add_paragraph_slide(prs, "Data Schemas & Models", analyzed.data_schemas)

    # 11) Target Users / Closing
    if analyzed.target_users:
        _add_paragraph_slide(prs, "Target Users & Use Cases", analyzed.target_users)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def render_ppt_from_analyzed(analyzed: "AnalyzedContent") -> bytes:
    """
    Phase 2b-PPT + Phase 3: AnalyzedContent → PPTX bytes.

    Builds a slide deck adapted to the repo type, then serializes.
    Requires ``python-pptx`` to be installed.
    """
    from pptx import Presentation

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = 9144000   # 10 inches
    prs.slide_height = 5143500  # 7.5 inches → 5.63 for 16:9

    _build_slides(prs, analyzed)

    buf = BytesIO()
    prs.save(buf)
    ppt_bytes = buf.getvalue()
    logger.info("PPT render complete — %d bytes, %d slides", len(ppt_bytes), len(prs.slides))
    return ppt_bytes


def render_ppt(summary_text: str, repo_name: str) -> bytes:
    """
    Backward-compatible entry point.

    Wraps the legacy summary_text in a minimal AnalyzedContent and
    delegates to the structured renderer.
    """
    from api.content_analyzer import AnalyzedContent

    analyzed = AnalyzedContent(repo_name=repo_name)
    # Parse the legacy summary_text into structured fields as best we can
    analyzed = _parse_legacy_summary(summary_text, repo_name)
    return render_ppt_from_analyzed(analyzed)


def _parse_legacy_summary(summary_text: str, repo_name: str) -> "AnalyzedContent":
    """
    Best-effort parse of the old 7-section plain-text summary into
    an AnalyzedContent so the structured PPT renderer can consume it.
    """
    from api.content_analyzer import AnalyzedContent, ModuleInfo, TechStack

    sections: dict[str, list[str]] = {}
    current_key = ""
    for line in summary_text.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        # Detect section headers (ends with ':')
        if stripped.endswith(":") and len(stripped) < 80 and not stripped.startswith("-"):
            current_key = stripped.rstrip(":")
            sections[current_key] = []
        elif current_key:
            sections[current_key].append(stripped)

    def bullets(key_prefix: str) -> list[str]:
        for k, v in sections.items():
            if k.lower().startswith(key_prefix.lower()):
                return [b.lstrip("- ").lstrip("* ") for b in v]
        return []

    def paragraph(key_prefix: str) -> str:
        for k, v in sections.items():
            if k.lower().startswith(key_prefix.lower()):
                return " ".join(v)
        return ""

    modules = []
    for b in bullets("key module"):
        if ":" in b:
            name, resp = b.split(":", 1)
            modules.append(ModuleInfo(name=name.strip(), responsibility=resp.strip()))
        else:
            modules.append(ModuleInfo(name=b, responsibility=""))

    return AnalyzedContent(
        repo_name=repo_name,
        project_overview=paragraph("project overview"),
        architecture=bullets("architecture"),
        tech_stack=TechStack(
            languages=[],
            frameworks=[],
            key_libraries=bullets("tech stack"),
            infrastructure=[],
        ),
        key_modules=modules,
        data_flow=bullets("data flow"),
        api_points=bullets("api"),
        target_users=paragraph("target user"),
    )
